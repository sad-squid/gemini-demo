from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2026 Gemini/AI Era Aesthetics
# Deep space/midnight blue background
BG_DARK = RGBColor(11, 15, 25)

# Ethereal/Nebula gradients (approximated with solid vibrant colors for python-pptx)
# For the AI "sparkle" we use sharp, luminous colors on dark background
GEMINI_BLUE = RGBColor(106, 176, 243)
GEMINI_PURPLE = RGBColor(190, 140, 255)
GEMINI_PINK = RGBColor(255, 138, 179)
GEMINI_CYAN = RGBColor(0, 230, 255)

TEXT_WHITE = RGBColor(248, 249, 250)
TEXT_MUTED = RGBColor(168, 176, 185)

# Glassmorphism/Translucent card approximations
CARD_BG = RGBColor(20, 26, 40) # Slightly lighter than background
CARD_BORDER = RGBColor(40, 48, 65)

def set_bg_dark(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_gemini_sparkle(slide):
    """Adds a subtle, modern 'Gemini' AI gradient spark effect (approximated) at the top."""
    # We will create overlapping very thin rectangles to fake a smooth glowing line
    width = prs.slide_width / 3
    height = Inches(0.05)
    
    # Left: Blue to Purple
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = GEMINI_BLUE
    shape1.line.fill.background()
    
    # Mid: Purple to Pink
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, width, 0, width, height)
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = GEMINI_PURPLE
    shape2.line.fill.background()
    
    # Right: Pink to Cyan
    shape3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, width*2, 0, width, height)
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = GEMINI_CYAN
    shape3.line.fill.background()

def add_title_slide(title, subtitle, author):
    blank_slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_dark(slide)
    add_gemini_sparkle(slide)
    
    # Subtle glowing orb effect behind title (faked with a large, low-opacity-looking circle)
    # python-pptx doesn't support transparency well, so we use a dark muted color
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.5), Inches(1), Inches(4.333), Inches(4.333))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(18, 24, 45) # Subtle glow
    glow.line.fill.background()
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Google Sans' # Will fallback to Arial if not installed on machine viewing it
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.333), Inches(1))
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.name = 'Google Sans'
    p2.font.size = Pt(28)
    p2.font.color.rgb = GEMINI_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(11.333), Inches(1))
    tf3 = author_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"Developer: {author}"
    p3.font.name = 'Google Sans'
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content, highlight_color=GEMINI_PURPLE):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_dark(slide)
    add_gemini_sparkle(slide)
    
    # Title Text (Clean, modern, left-aligned)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
    tf_title = title_box.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Google Sans'
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.alignment = PP_ALIGN.LEFT
    
    # A subtle glowing dot next to the title
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.85), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = highlight_color
    dot.line.fill.background()

    # Glassmorphic Content Card
    content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(11.333), Inches(5))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = CARD_BG
    content_box.line.color.rgb = CARD_BORDER
    content_box.line.width = Pt(1)
    
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.5)
    tf2.margin_right = Inches(0.5)
    tf2.margin_top = Inches(0.5)
    
    p0 = tf2.paragraphs[0]
    p0.text = content[0]
    p0.font.name = 'Google Sans'
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = highlight_color
    
    for item in content[1:]:
        p = tf2.add_paragraph()
        p.text = item
        p.font.name = 'Google Sans'
        if item.startswith('-') or item.startswith('*'):
            p.level = 1
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_MUTED
            p.space_before = Pt(12)
        else:
            p.level = 0
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.space_before = Pt(24)

# Slide 1
add_title_slide("Local Lens", "Hyperlocalized Context Builder\nPowered by Gemini Advanced", "Cuong Duong (cuongduong.git@gmail.com)")

# Slide 2
add_content_slide("1. The Problem", [
    "Physical Discoveries Fade Quickly",
    "- You walk down the street and see a flyer for an underground concert or event.",
    "- You snap a picture, but it just sits in your camera roll.",
    "- There is no instant, frictionless way to convert visual, real-world objects into mapped, actionable, and structured data."
], GEMINI_PINK)

# Slide 3
add_content_slide("2. The Solution: Local Lens", [
    "From Raw Pixel to Enriched Intelligence",
    "- Local Lens instantly bridges the physical world with a digital, interactive hyperlocal map.",
    "- Upload a single photo: Local Lens uses Multimodal Gemini AI to generate a structured model.",
    "- Agentic Search fetches the missing pieces (coordinates, links, schedules) autonomously."
], GEMINI_CYAN)

# Slide 4
add_content_slide("3. How It Works", [
    "The 3-Step Magic",
    "- 1. Capture & Parse: User uploads a photo. Gemini Vision API instantly extracts unstructured text to precise JSON.",
    "- 2. Agentic Enrichment: Autonomous agent queries the web to find exact geo-coordinates and ticketing links.",
    "- 3. Interactive Map: Data drops onto a live map. Users ask the Local Lens Chatbot for recommendations."
], GEMINI_BLUE)

# Slide 5
add_content_slide("4. Architecture & Tech Stack", [
    "Powered by Google Cloud",
    "- Frontend: React/Next.js with interactive Mapbox UI.",
    "- Backend: Node.js/Express API orchestration.",
    "- AI & Agents: Gemini Multimodal Vision API & Gemini Text API (Chatbot).",
    "- Database: Google Cloud Firestore (Venues, events, geo-indices).",
    "- Infrastructure: Google Cloud Run."
], GEMINI_PURPLE)

# Slide 6
add_content_slide("5. Innovation", [
    "Leveraging the Latest Gemini Capabilities",
    "- Multimodal Generation: The AI sees the flyer, understands visual hierarchy, and extracts meaning automatically.",
    "- Managed Agent Workflows: Background agents actively crawl the web to find missing context.",
    "- Grounded Chat: Local Lens concierge grounds its answers in entities stored in our Firestore."
], GEMINI_PINK)

# Slide 7
add_content_slide("6. Core Features", [
    "What you can do:",
    "- Snap & Build: Frictionless visual ingestion of data.",
    "- Pulse Feed & Vibe Filters: Explore Tokyo events by vibe (e.g., Cyberpunk, Cozy, Underground).",
    "- Local Lens Concierge: Ask 'What acoustic gigs are near Shibuya tonight?'",
    "- Real-time Shared Map: A collaborative map built entirely from user-submitted photography."
], GEMINI_CYAN)

# Slide 8
add_content_slide("7. Demo Walkthrough", [
    "What You'll See in the Video",
    "- Uploading a raw, complex event flyer image.",
    "- UI displaying the successfully structured JSON extraction.",
    "- Background agent enriching location into latitude/longitude coordinates.",
    "- Event dynamically popping up on our interactive map.",
    "- Asking Local Lens Chatbot to plan an itinerary."
], GEMINI_BLUE)

# Slide 9
add_content_slide("8. Thank You!", [
    "Ready to explore the real world, better.",
    "- Developer: Cuong Duong",
    "- Source Code: github.com/sad-squid/gemini-demo",
    "- Live Portal: Local Lens App",
    "- (All hackathon requirements—Google Cloud usage, Innovation, and Completeness—have been met!)"
], GEMINI_PURPLE)

prs.save("Local_Lens_Pitch_Deck.pptx")
